import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Brand Colors
    BG_DARK = RGBColor(7, 31, 36)      # #071F24
    BG_CARD = RGBColor(11, 43, 51)     # #0B2B33
    TEAL = RGBColor(2, 128, 144)       # #028090
    MINT = RGBColor(2, 195, 154)       # #02C39A
    CYAN = RGBColor(56, 189, 248)      # #38BDF8
    WHITE = RGBColor(255, 255, 255)
    SLATE_LIGHT = RGBColor(203, 213, 225) # #CBD5E1
    SLATE_MUTED = RGBColor(148, 163, 184) # #94A3B8
    ROSE = RGBColor(244, 63, 94)       # #F43F5E
    AMBER = RGBColor(251, 191, 36)     # #FBBF24

    blank_slide_layout = prs.slide_layouts[6]

    def add_blank_slide_with_bg():
        slide = prs.slides.add_slide(blank_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_DARK
        return slide

    def add_header(slide, tag_text, title_text, tag_color=MINT):
        # Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf_tag = tag_box.text_frame
        tf_tag.word_wrap = True
        p_tag = tf_tag.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = tag_color

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

    def add_card(slide, left, top, width, height, title, desc, title_color=MINT):
        # Card Background Shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_CARD
        shape.line.color.rgb = TEAL
        shape.line.width = Pt(1)

        # Content inside card
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = title_color
        p1.space_after = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = SLATE_LIGHT

    # =========================================================================
    # SLIDE 1: Cover & National Vision
    # =========================================================================
    s1 = add_blank_slide_with_bg()
    
    logo_path = r"c:\Users\Doston\Downloads\SAPAR\sapar-logo-512.png"
    if os.path.exists(logo_path):
        s1.shapes.add_picture(logo_path, Inches(5.8), Inches(0.6), Inches(1.7), Inches(1.7))

    badge_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.4))
    p_badge = badge_box.text_frame.paragraphs[0]
    p_badge.text = "★ PRESIDENT TECH AWARD 2026 NOMZODI"
    p_badge.font.size = Pt(12)
    p_badge.font.bold = True
    p_badge.font.color.rgb = MINT
    p_badge.alignment = PP_ALIGN.CENTER

    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.8))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "Oʻzbekiston Biznesini Boyitish va Milliy Iqtisodiyotni Yuksaltirish Uchun Yagona Bulutli ERP"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "Buxgalteriya (21-BHMS) • Soliq & Didox • FIFO Ombor • Savdo & TTN • Sensorli POS • HRM • Bank Skoringi • CRM"
    p2.font.size = Pt(13)
    p2.font.bold = False
    p2.font.color.rgb = CYAN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    p3 = tf1.add_paragraph()
    p3.text = "Buxgalteriya hisobi (21-BHMS), Soliq integratsiyasi, FIFO ombor, Savdo va TTN, Sensorli POS Kassa, HRM & Ish haqi, Bank skoringi va CRM."
    p3.font.size = Pt(13)
    p3.font.color.rgb = SLATE_LIGHT
    p3.space_before = Pt(8)

    add_card(s1, Inches(0.8), Inches(4.3), Inches(2.7), Inches(2.2), "100% Bulutli", "Oʻrnatish talab etilmaydi, barcha brauzer va mobil qurilmalardan ishlaydi.")
    add_card(s1, Inches(3.8), Inches(4.3), Inches(2.7), Inches(2.2), "Bank Skoringi", "Tadbirkorlarga 1 kunda garov taʼminotisiz aylanma mablagʻ kreditlari.")
    add_card(s1, Inches(6.8), Inches(4.3), Inches(2.7), Inches(2.2), "E-IMZO & Soliq", "Didox, Factura.uz va Soliq qoʻmitasi bilan toʻliq avtomatlashtirilgan.")
    add_card(s1, Inches(9.8), Inches(4.3), Inches(2.7), Inches(2.2), "5x Hamyonbop", "1C tizimiga nisbatan har bir biznesga $8,500+ mablagʻni tejash imkoniyati.")

    # =========================================================================
    # SLIDE 2: Macroeconomic Challenges
    # =========================================================================
    s2 = add_blank_slide_with_bg()
    add_header(s2, "Makroiqtisodiy Muammolar", "Oʻzbekiston Biznesi va Milliy Iqtisodiyotning Asosiy Toʻsiqlari", ROSE)
    add_card(s2, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "$120M+ Valyutaning Xorijga Chiqishi", "Oʻzbekiston korxonalari har yili xorijiy eskirgan dasturlarga (1C litsenziyalari, serverlar, qimmat maslahatchilar) millionlab dollar xorijiy valyuta sarflamoqda.", ROSE)
    add_card(s2, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "Bank Krediti Olishdagi Qiyinchiliklar", "Kichik biznesning buxgalteriyasi va real aylanmasi shaffof koʻrinmagani sababli, banklar ularga tezkor va garovsiz aylanma mablagʻ kreditlarini bera olmaydi.", AMBER)
    add_card(s2, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "Tarqoq Dasturlar va Xufyona Iqtisodiyot", "POS kassa alohida, ombor daftarda, fakturalar Didoxda, hisoblar bankda. Maʼlumotlar tarqoqligi korxonaning daromad yoʻqotishiga sabab boʻlmoqda.", CYAN)

    # =========================================================================
    # SLIDE 3: National Mission
    # =========================================================================
    s3 = add_blank_slide_with_bg()
    add_header(s3, "Milliy Missiya", "Bizning Maqsadimiz: Bizneslarni Kuchli va Boy Qilish", MINT)
    add_card(s3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "+25% Sof Daromadni Oshirish", "FIFO tannarxini aniq hisoblash, tovar kamomadini yoʻqotish va savdo quvuri orqali korxona sof rentabelligi 25% gacha oshadi.", MINT)
    add_card(s3, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "5x Xarajatlarni Qisqartirish", "Qimmat jismoniy serverlar, doimiy 1C dasturchilari va texnik xizmat koʻrsatish xarajatlari toʻliq tejaladi.", CYAN)
    add_card(s3, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "1 Kunda Bank Krediti", "SAPAR dagi tasdiqlangan moliyaviy hisobotlar orqali banklar tadbirkorga 1 kunda garovsiz va qulay foizli aylanma mablagʻ kreditlarini ajratadi.", WHITE)

    # =========================================================================
    # SLIDE 4: Competitor Matrix Table
    # =========================================================================
    s4 = add_blank_slide_with_bg()
    add_header(s4, "Raqobat Tahlili", "Nima Uchun SAPAR ERP Bozor Yechimlaridan Ustun?", MINT)
    
    rows, cols = 6, 5
    table_shape = s4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    table = table_shape.table

    headers = ["Xususiyat / Feature", "🚀 SAPAR ERP", "🏛️ 1C:Enterprise", "🛒 Poster POS", "🌍 Odoo / Microsoft"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL if col_idx != 1 else MINT

    data = [
        ["100% Bulutli va Mobil", "✅ Oʻrnatish shart emas", "❌ Server va oʻrnatish shart", "⚠️ Faqat POS qismi", "✅ Bulutli"],
        ["21-BHMS & 1/2-Shakl Balans", "✅ Nativ / Tayyor oʻrnatilgan", "✅ Mavjud", "❌ Mavjud emas", "❌ Maxsus sozlash talab etiladi"],
        ["E-IMZO, Didox & Soliq API", "✅ Toʻgʻridan-toʻgʻri integratsiya", "⚠️ Qimmat qoʻshimcha modul", "❌ Mavjud emas", "❌ Mavjud emas"],
        ["Bank Skoringi va Kredit API", "✅ Oʻrnatilgan moliyaviy skoring", "❌ Mavjud emas", "❌ Mavjud emas", "❌ Mavjud emas"],
        ["5 Yillik Xarajat (TCO)", "💰 $1,200 – $1,800", "💸 $9,700+", "💰 $3,500+", "💸 $20,000+"]
    ]

    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_CARD

    # =========================================================================
    # SLIDE 5: 5-Year TCO Breakdown
    # =========================================================================
    s5 = add_blank_slide_with_bg()
    add_header(s5, "Xarajatlar Tahlili", "5 Yillik Egalik Qilish Xarajati: 1C ga Nisbatan $8,500+ Tejash", MINT)
    add_card(s5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.7), "🏛️ 1C:Enterprise (5 Yillik Xarajat)", "• Dastlabki litsenziyalar va konfiguratsiyalar: $1,200\n• Jismoniy server va UPS xaridi: $1,500\n• Oylik 1C mutaxassisi xizmati ($100/oy x 60 oy): $6,000\n• Yillik yangilanishlar va qonunchilik moslashuvi: $1,000\n\nJAMI 5 YILLIK XARAJAT: $9,700+", ROSE)
    add_card(s5, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.7), "🚀 SAPAR Cloud ERP (5 Yillik Xarajat)", "• Server va bulut infratuzilmasi: $0 (Kiritilgan)\n• Dasturchi yollash va texnik xizmat: $0 (Avtomatik yangilanadi)\n• 5 Yillik toʻliq SaaS obuna paketi: $1,500\n• 24/7 Doimiy texnik qoʻllab-quvvatlash: Kiritilgan\n\nHAR BIR BIZNESGA SOF TEJAMKORLIK: +$8,200 TEJALDI!", MINT)

    # =========================================================================
    # SLIDE 6: Banking Partnerships
    # =========================================================================
    s6 = add_blank_slide_with_bg()
    add_header(s6, "Banklar Bilan Hamkorlik", "Bank Integratsiyasi va Tadbirkorlarga Tezkor Garovsiz Kredit", CYAN)
    add_card(s6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "Moliyaviy Skoring Dvigateli", "Hamkor banklar (Ipak Yoʻli, Kapitalbank, SQB, Agrobank) SAPAR dagi shaffof savdo aylanmasi asosida tadbirkorga 1 kunda kredit liniyalarini ajratadi.", CYAN)
    add_card(s6, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "1C:ClientBank Avto-Sinxronizatsiya", "Oʻzbekistonning barcha banklaridan koʻchirmalar (vipiska) 1 soniyada yuklanadi va buxgalteriya schotlari bilan avtomatik solishtiriladi.", MINT)
    add_card(s6, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "Payme, Click va Uzum Pay", "Mijozlardan QR-kod orqali toʻlovlarni toʻgʻridan-toʻgʻri korxona hisob-raqamiga qabul qilish va darhol buxgalteriyada aks ettirish.", WHITE)

    # =========================================================================
    # SLIDE 7: Two-Tier Navigation vs Existing Products
    # =========================================================================
    s7 = add_blank_slide_with_bg()
    add_header(s7, "Dizayn va UX", "Mavjud Mahsulotlarga Nisbatan 2-Bosqichli Tezkor Navigatsiya", MINT)
    add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.7), "Fokuslangan 2-Bosqichli Struktura", "• 1-Bosqich (72px Rail): 13 ta asosiy boʻlimga 1 bosishda oʻtish.\n• 2-Bosqich (256px Panel): Tezkor qidiruv, + Yangi... tugmalari va filtrlar.\n• 1-Bosishda Mahkamlash (Pin): Katta moliyaviy jadvallar uchun 100% keng ekran.\n• Dual Sidebar Arxitekturasi: Sahifa qayta yuklanmasdan, darhol kerakli ish stoli ochiladi.\n\nMavjud eskirgan mahsulotlarning 80+ menyudan iborat chalkash roʻyxatidan xalos boʻling.", MINT)
    
    dash_img = r"c:\Users\Doston\Downloads\SAPAR\dashboard_ui.png"
    if os.path.exists(dash_img):
        s7.shapes.add_picture(dash_img, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.7))
    else:
        add_card(s7, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.7), "Dual Sidebar Arxitekturasi", "Foydalanuvchi qaysi sahifada boʻlsa, tizim oʻsha boʻlimni avtomatik aniqlaydi.", CYAN)

    # =========================================================================
    # SLIDE 8: 21-BHMS Accounting Standards
    # =========================================================================
    s8 = add_blank_slide_with_bg()
    add_header(s8, "Buxgalteriya va Moliya", "21-son BHMS Milliy Buxgalteriya Standartlari", MINT)
    add_card(s8, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "Schotlar Rejasi va Provodkalar", "0100 dan 9900 gacha barcha milliy buxgalteriya schotlari tayyor kiritilgan. Har bir amaliyot avtomatik Bosh kitobga yoziladi.", MINT)
    add_card(s8, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "1-Shakl, 2-Shakl va Oborotka", "Buxgalteriya balansi (1-shakl), Moliyaviy natijalar (2-shakl / P&L) va Aylanma vedomost (Oborotka) 1 soniyada avtomatik shakllanadi.", CYAN)
    add_card(s8, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "Vzaimozachet va Koʻp Valyutali Hisob", "Korxonalar oʻrtasidagi oʻzaro qarzdorliklarni yopish (Contras), UZS, USD, EUR, RUB kurslari boʻyicha avtomatik qayta baholash.", WHITE)

    # =========================================================================
    # SLIDE 9: Soliq, Didox & E-IMZO
    # =========================================================================
    s9 = add_blank_slide_with_bg()
    add_header(s9, "Soliq va Raqamli Imzo", "E-IMZO, Didox va Soliq Qoʻmitasi Bilan Toʻliq Integratsiya", MINT)
    add_card(s9, Inches(0.8), Inches(1.8), Inches(2.7), Inches(4.7), "Nativ E-IMZO Imzolash", "USB e-token yoki .pfx sertifikat orqali brauzerdan chiqmasdan toʻgʻridan-toʻgʻri PKCS#7 formatida hujjatlarni imzolash.", MINT)
    add_card(s9, Inches(3.8), Inches(1.8), Inches(2.7), Inches(4.7), "Didox va Factura.uz", "Elektron fakturalarni 1 bosishda yuborish, statusini kuzatish va kiruvchi hujjatlarni avtomatik qabul qilish.", CYAN)
    add_card(s9, Inches(6.8), Inches(1.8), Inches(2.7), Inches(4.7), "MXIK / IKPU Kodlari", "Tovar va xizmatlar uchun milliy klassifikator kodlari va shtrix-kodlarni avtomatik toʻgʻri belgilash.", WHITE)
    add_card(s9, Inches(9.8), Inches(1.8), Inches(2.7), Inches(4.7), "Soliq Deklaratsiyalari", "QQS (Forma 10006_29), JShODS (11101_14) va Aylanmadan olinadigan soliq (10104_18) hisobotlari.", MINT)

    # =========================================================================
    # SLIDE 10: Sales & TTN
    # =========================================================================
    s10 = add_blank_slide_with_bg()
    add_header(s10, "Savdo Zanjiri", "Savdo Boshqaruvi, Hisob-Fakturalar va Yuk Xatlari (TTN)", MINT)
    add_card(s10, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "Tijorat Takliflari (KP)", "Mijozlarga rasmiy shartlar koʻrsatilgan taklifnomalarni yuborish va mijoz qabul qilgach, 1 bosishda toʻlov fakturasiga aylantirish.", MINT)
    add_card(s10, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "Yuk Xatlari (TTN)", "Tovarni ombordan mijozga joʻnatishda avtotransport vositasi va haydovchi maʼlumotlari qayd etiladigan milliy yuk xatlari.", CYAN)
    add_card(s10, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "Ommaviy Havolalar va QR-Kod", "Mijozga Telegram orqali hisob-fakturaning onlayn havolasini yuborish va QR-kod orqali toʻlovlarni qabul qilish.", WHITE)

    # =========================================================================
    # SLIDE 11: Multi-Warehouse & FIFO
    # =========================================================================
    s11 = add_blank_slide_with_bg()
    add_header(s11, "Ombor va Tovar Nazorati", "Koʻp Omborli Boshqaruv va FIFO Tannarx Qatlamlari", CYAN)
    add_card(s11, Inches(0.8), Inches(1.8), Inches(2.7), Inches(4.7), "FIFO Tannarx Hisobi", "Dastlabki kirim partiyalari narxi boʻyicha tovar tannarxi va sof foydani aniq hisoblash.", MINT)
    add_card(s11, Inches(3.8), Inches(1.8), Inches(2.7), Inches(4.7), "Filiallar va Omborlar", "Barcha filial va omborlardagi tovar qoldiqlarini real vaqt rejimida koʻrish.", CYAN)
    add_card(s11, Inches(6.8), Inches(1.8), Inches(2.7), Inches(4.7), "Omborlararo Koʻchirish", "Omborlar oʻrtasida tovarlarni xavfsiz va hujjatlashtirilgan holda oʻtkazish.", WHITE)
    add_card(s11, Inches(9.8), Inches(1.8), Inches(2.7), Inches(4.7), "Inventarizatsiya", "Ortiqcha tovar kirimi va kamomadni hisobdan chiqarish dalolatnomalari.", ROSE)

    # =========================================================================
    # SLIDE 12: Touch POS Terminal
    # =========================================================================
    s12 = add_blank_slide_with_bg()
    add_header(s12, "Chakana Savdo va POS", "Sensorli Touch POS Terminali va Smenalar Boshqaruvi", CYAN)
    add_card(s12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.7), "Kassa Operatsiyalari", "• Tezkor Sensorli Interfeys: Planshet va monitorlar uchun moslashtirilgan sub-soniyalik tezkor kassa oynasi.\n• Aralash (Split) Toʻlovlar: Naqd pul + Uzcard/Humo + Payme/Click + Nasiya (qarz) hisobiga boʻlib toʻlash.\n• Kassir Smenalari & X/Z Hisobotlar: Smenani ochish, kassa naqd pulini nazorat qilish va kun yakunida avtomatik X/Z hisoboti.\n• Fiskal cheklar va offline ishlash imkoniyati.", CYAN)
    
    pos_img = r"c:\Users\Doston\Downloads\SAPAR\pos_terminal.png"
    if os.path.exists(pos_img):
        s12.shapes.add_picture(pos_img, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.7))
    else:
        add_card(s12, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.7), "Kassir Smenalari va X/Z Hisobotlar", "Smenani ochish, kassa naqd pulini nazorat qilish va kun yakunida avtomatik X/Z hisoboti.", WHITE)

    # =========================================================================
    # SLIDE 13: HRM & Payroll
    # =========================================================================
    s13 = add_blank_slide_with_bg()
    add_header(s13, "HRM va Xodimlar", "Ish Vaqti Tabeli va Oʻzbekiston Oylik Ish Haqi Dvigateli", ROSE)
    add_card(s13, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "Ish Vaqti Tabeli (Tabel)", "Xodimlarning keldi-ketdi soatlari, dam olish kunlari va taʼtillarini hisobga oluvchi qonuniy tabel.", MINT)
    add_card(s13, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "12% + 12% + 0.1% Soliqlar", "JShODS (12%), Ijtimoiy soliq (12% / IT Park 1%) va INPS (0.1%) pensiya badallari avtomatik hisoblanadi.", ROSE)
    add_card(s13, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "Taʼtillar va Kompensatsiyalar", "Mehnat taʼtillari kunlari hisobi, kasallik varaqalari va kompensatsiya toʻlovlari.", WHITE)

    # =========================================================================
    # SLIDE 14: CRM Deals Pipeline
    # =========================================================================
    s14 = add_blank_slide_with_bg()
    add_header(s14, "CRM va Savdo Oʻsishi", "Vizual Bitimlar Quvuri va 360° Mijoz Balansi", MINT)
    add_card(s14, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.7), "Mijozlar va Savdo Quvuri", "• Kanban Bitimlar Quvuri: Lidlar ➔ Birinchi aloqa ➔ Tijorat taklifi ➔ Shartnoma ➔ Toʻlov bosqichlarini vizual boshqaring.\n• 360° Mijoz Kartasi: Mijozning barcha hisob-fakturalari, oʻzaro qarzdorlik balansi va toʻlovlari yagona oynada.\n• Telegram & WhatsApp Xabarlari: Mijozlarga hisob-faktura va toʻlov eslatmalarini avtomatik yuborish.\n• Mijozlar bilan ishlash tezligi 2 barobar oshadi.", MINT)
    
    crm_img = r"c:\Users\Doston\Downloads\SAPAR\crm_kanban.png"
    if os.path.exists(crm_img):
        s14.shapes.add_picture(crm_img, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.7))
    else:
        add_card(s14, Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.7), "Telegram & WhatsApp Xabarlari", "Mijozlarga hisob-faktura va toʻlov eslatmalarini avtomatik yuborish.", WHITE)

    # =========================================================================
    # SLIDE 15: SAPAR AI Co-Pilot
    # =========================================================================
    s15 = add_blank_slide_with_bg()
    add_header(s15, "Sunʼiy Intellekt", "SAPAR AI Co-Pilot — Aqlli Biznes Maslahatchisi", CYAN)
    add_card(s15, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "Cheklarni Tanib Olish (OCR)", "Xarid cheklarini rasmga oling — AI uni avtomatik tovar kirimi va xarajatga aylantiradi.", MINT)
    add_card(s15, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "Moliyaviy Bashorat va Maslahat", "Kutilayotgan soliq toʻlovlari va pul tanqisligi (cash gap) boʻyicha rahbarga aqlli tavsiyalar.", CYAN)
    add_card(s15, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "Tabiiy Tilda Savol-Javob", "\"Bu oy sof foydamiz qancha boʻldi?\" deb soʻrang — AI hisobotni darhol chiqaradi.", WHITE)

    # =========================================================================
    # SLIDE 16: Macroeconomic Impact
    # =========================================================================
    s16 = add_blank_slide_with_bg()
    add_header(s16, "Iqtisodiy Samara", "Oʻzbekiston Iqtisodiyoti va Xalq Farovonligiga Milliy Taʼsiri", MINT)
    add_card(s16, Inches(0.8), Inches(1.8), Inches(2.7), Inches(4.7), "$120M+ Tejash", "Import Oʻrnini Bosish: Xorijiy dasturlarga ketayotgan millionlab dollar valyutani mamlakatimizda saqlab qolish.", MINT)
    add_card(s16, Inches(3.8), Inches(1.8), Inches(2.7), Inches(4.7), "+35% Unumdorlik", "50,000+ korxonada qogʻozbozlikni yoʻqotish va operatsion tezlikni oshirish.", CYAN)
    add_card(s16, Inches(6.8), Inches(1.8), Inches(2.7), Inches(4.7), "100% Shaffoflik", "Xufyona aylanmani qonuniylashtirish va davlat byudjetiga tushumlarni koʻpaytirish.", WHITE)
    add_card(s16, Inches(9.8), Inches(1.8), Inches(2.7), Inches(4.7), "Markaziy Osiyo", "Mahalliy korxonalarni mintaqa bozorlariga kengaytirish.", MINT)

    # =========================================================================
    # SLIDE 17: 3-Year Financial Model
    # =========================================================================
    s17 = add_blank_slide_with_bg()
    add_header(s17, "Moliyaviy Model", "3 Yillik Moliyaviy Reja va SaaS Unit Economics", MINT)
    add_card(s17, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "2026 (1-YIL): $750,000", "• 2,500 Faol Korxona\n• Oylik daromad: $62,500\n• Sof Foyda Rentabelligi: 32%\n• CAC: $45 | LTV: $950+", MINT)
    add_card(s17, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "2027 (2-YIL): $3,600,000", "• 12,000 Faol Korxona\n• Oylik daromad: $300,000\n• Sof Foyda Rentabelligi: 48%\n• LTV / CAC nisbati: 21x", CYAN)
    add_card(s17, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "2028 (3-YIL): $15,000,000", "• 50,000 Faol Korxona\n• Oylik daromad: $1,250,000\n• Sof Foyda Rentabelligi: 56%\n• Churn darajasi: < 1.2% / oy", WHITE)

    # =========================================================================
    # SLIDE 18: Security & Compliance
    # =========================================================================
    s18 = add_blank_slide_with_bg()
    add_header(s18, "Xavfsizlik va Arxitektura", "Bank Darajasidagi Xavfsizlik va OʻRQ-547 Standarti", MINT)
    add_card(s18, Inches(0.8), Inches(1.8), Inches(2.7), Inches(4.7), "PostgreSQL & Prisma", "Millionlab tranzaksiyalarga bardosh beruvchi relyatsion maʼlumotlar bazasi.", MINT)
    add_card(s18, Inches(3.8), Inches(1.8), Inches(2.7), Inches(4.7), "Multi-Tenancy", "Har bir korxona uchun alohida subdomen (korxona.sapar.uz) va toʻliq maʼlumotlar izolyatsiyasi.", CYAN)
    add_card(s18, Inches(6.8), Inches(1.8), Inches(2.7), Inches(4.7), "AES-256 Shifrlash", "Parollar va maʼlumotlarni bank darajasida xavfsiz shifrlab saqlash.", WHITE)
    add_card(s18, Inches(9.8), Inches(1.8), Inches(2.7), Inches(4.7), "Oʻzbekiston OʻRQ-547", "Barcha maʼlumotlar Oʻzbekiston hududidagi maʼlumotlar markazida saqlanadi.", MINT)

    # =========================================================================
    # SLIDE 19: Pricing & Deployment
    # =========================================================================
    s19 = add_blank_slide_with_bg()
    add_header(s19, "Tariflar va Oʻrnatish", "Hamyonbop va Moslashuvchan Biznes Modeli", MINT)
    add_card(s19, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.7), "STARTUP ($15 / oy)", "• 3 tagacha foydalanuvchi\n• 1 ta ombor va Hisob-fakturalar\n• Didox va E-IMZO integratsiyasi\n• Kichik doʻkonlar va savdo korxonalari uchun.", MINT)
    add_card(s19, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.7), "BUSINESS ($35 / oy)", "• 10 tagacha foydalanuvchi\n• 21-BHMS Buxgalteriya va 1/2-Shakl\n• Sensorli POS Kassa va Smenalar\n• Oʻrta biznes va ulgurji savdo korxonalari.", CYAN)
    add_card(s19, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.7), "ENTERPRISE ($80 / oy)", "• Cheksiz foydalanuvchilar\n• On-Premises Docker oʻrnatish\n• 24/7 Shaxsiy qoʻllab-quvvatlash\n• Yirik korxonalar va maxsus serverlar.", WHITE)

    # =========================================================================
    # SLIDE 20: Leadership & Engineering Team
    # =========================================================================
    s20 = add_blank_slide_with_bg()
    add_header(s20, "Bizning Jamoa", "Kuchli va Tajribali Mutaxassislar Jamoasi", MINT)
    add_card(s20, Inches(0.8), Inches(1.8), Inches(3.6), Inches(2.2), "Dostonbek & Hammuassislar", "Bosh Ijrochi Direktor (CEO & Founder)\nSaaS strategiyasi, xalqaro tajriba, korporativ boshqaruv va ekotizim viziyasi.", MINT)
    add_card(s20, Inches(4.8), Inches(1.8), Inches(3.6), Inches(2.2), "Bosh Texnologik Arxitektor", "Bosh Texnik Direktor (CTO & Lead Architect)\nYuqori yuklamali PostgreSQL, Node.js, microservices va bank darajasidagi xavfsizlik.", CYAN)
    add_card(s20, Inches(8.8), Inches(1.8), Inches(3.6), Inches(2.2), "Bosh Auditor & Moliyaviy Ekspert", "Bosh Moliyachi (CFO & Tax Lead)\n21-son BHMS, Soliq kodeksi, 1/2-shakllar va milliy soliq deklaratsiyalari mutaxassisi.", WHITE)
    add_card(s20, Inches(0.8), Inches(4.3), Inches(3.6), Inches(2.2), "Fintech & Bank Integratsiyasi", "Fintech & Bank API Lead\nDidox, E-IMZO, 1C:ClientBank va tijorat banklari skoring tizimlari boʻyicha muhandis.", CYAN)
    add_card(s20, Inches(4.8), Inches(4.3), Inches(3.6), Inches(2.2), "Sunʼiy Intellekt Muhandisi", "Head of AI & OCR Extraction\nCheklarni rasmga olib tanish (OCR), moliyaviy tahlil LLM modellari va algoritmik tahlil.", ROSE)
    add_card(s20, Inches(8.8), Inches(4.3), Inches(3.6), Inches(2.2), "Joriy Qilish & Mijozlar Xizmati", "Head of Customer Success & Support\n500+ korxonalarda ERP joriy qilish, xodimlarni oʻqitish va 24/7 tezkor texnik yordam.", MINT)

    # =========================================================================
    # SLIDE 21: Summary, Contact & Call-to-Action
    # =========================================================================
    s21 = add_blank_slide_with_bg()
    
    if os.path.exists(logo_path):
        s21.shapes.add_picture(logo_path, Inches(6.0), Inches(0.4), Inches(1.3), Inches(1.3))

    title_box21 = s21.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.4))
    tf21 = title_box21.text_frame
    tf21.word_wrap = True
    p21_1 = tf21.paragraphs[0]
    p21_1.text = "Biznesingizni SAPAR Bilan Boshqaring!"
    p21_1.font.size = Pt(32)
    p21_1.font.bold = True
    p21_1.font.color.rgb = WHITE
    p21_1.alignment = PP_ALIGN.CENTER

    p21_2 = tf21.add_paragraph()
    p21_2.text = "Eskirgan va qimmat dasturlardan xalos boʻling. Yagona oynada barcha korxona jarayonlarini avtomatlashtiring."
    p21_2.font.size = Pt(14)
    p21_2.font.color.rgb = CYAN
    p21_2.alignment = PP_ALIGN.CENTER
    p21_2.space_before = Pt(6)

    add_card(s21, Inches(0.8), Inches(3.2), Inches(3.6), Inches(2.2), "Rasmiy Veb-Sayt", "https://sapar.uz\nOnlayn roʻyxatdan oʻtish va tariflar.", MINT)
    add_card(s21, Inches(4.8), Inches(3.2), Inches(3.6), Inches(2.2), "Elektron Pochta", "info@sapar.uz\nKorporativ soʻrovlar va hamkorlik.", CYAN)
    add_card(s21, Inches(8.8), Inches(3.2), Inches(3.6), Inches(2.2), "Manzil & Aloqa", "Toshkent shahri, IT Park Oʻzbekiston\nJonli Demo: http://localhost:3000/admin", WHITE)

    # CTA Button / Banner
    cta_shape = s21.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(5.8), Inches(5.7), Inches(0.9))
    cta_shape.fill.solid()
    cta_shape.fill.fore_color.rgb = MINT
    cta_tf = cta_shape.text_frame
    cta_p = cta_tf.paragraphs[0]
    cta_p.text = "DEMONI SINAB KOʻRISH (http://localhost:3000/admin)"
    cta_p.font.size = Pt(14)
    cta_p.font.bold = True
    cta_p.font.color.rgb = BG_DARK
    cta_p.alignment = PP_ALIGN.CENTER

    # Save
    output_filename = "SAPAR_ERP_Enterprise_Presentation_2026.pptx"
    prs.save(output_filename)
    print(f"Successfully generated 21-slide presentation: {output_filename}")

if __name__ == "__main__":
    create_presentation()
